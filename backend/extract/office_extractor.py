from pathlib import Path
from typing import List

from backend.config import MAX_CONTENT_LENGTH
from backend.extract.base import ExtractResult
from backend.logger import logger


def _iter_docx_tables(doc) -> List[str]:
    cells: List[str] = []
    try:
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        cells.append(t)
    except Exception:
        pass
    return cells


def extract_docx(path: str | Path) -> ExtractResult:
    p = Path(path)
    try:
        from docx import Document
    except ImportError:
        return ExtractResult(success=False, error="python-docx not installed")

    try:
        doc = Document(str(p))
    except Exception as e:
        logger.warning(f"open docx failed {p}: {e}")
        return ExtractResult(success=False, error=f"open docx: {e}")

    parts: List[str] = []
    try:
        for para in doc.paragraphs:
            t = para.text
            if t:
                parts.append(t)
    except Exception as e:
        logger.warning(f"docx paragraphs error {p}: {e}")

    parts.extend(_iter_docx_tables(doc))

    text = "\n".join(parts)
    if len(text) > MAX_CONTENT_LENGTH:
        text = text[:MAX_CONTENT_LENGTH]

    return ExtractResult(success=True, content=text)


def extract_pptx(path: str | Path) -> ExtractResult:
    p = Path(path)
    try:
        from pptx import Presentation
    except ImportError:
        return ExtractResult(success=False, error="python-pptx not installed")

    try:
        prs = Presentation(str(p))
    except Exception as e:
        logger.warning(f"open pptx failed {p}: {e}")
        return ExtractResult(success=False, error=f"open pptx: {e}")

    parts: List[str] = []
    slide_num = 0
    try:
        for slide in prs.slides:
            slide_num += 1
            slide_texts: List[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text:
                            slide_texts.append(run.text)
            if slide_texts:
                parts.append(f"=== Slide {slide_num} ===\n" + "\n".join(slide_texts))
    except Exception as e:
        logger.warning(f"pptx extract error {p}: {e}")

    text = "\n\n".join(parts)
    if len(text) > MAX_CONTENT_LENGTH:
        text = text[:MAX_CONTENT_LENGTH]

    return ExtractResult(success=True, content=text, metadata={"slides": slide_num})


def extract_xlsx(path: str | Path) -> ExtractResult:
    return ExtractResult(
        success=True,
        content="",
        metadata={"note": "xlsx text extraction skipped (use csv or install openpyxl)"},
    )
